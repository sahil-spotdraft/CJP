#!/usr/bin/env python3
"""Generate the CJP Moonshot Feature Hub hackathon pitch deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Design system ---
BG = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
EMERALD_DIM = RGBColor(0x06, 0x78, 0x5A)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).resolve().parents[1] / "pitch" / "CJP-Moonshot-Hackathon-Pitch.pptx"


def _set_run(run, text: str, size: int, color: RGBColor, bold: bool = False) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    *,
    paragraphs: list[tuple[str, int, RGBColor, bool]] | None = None,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
):
    """Add a text box with one or more styled paragraphs.

    paragraphs: list of (text, size_pt, color, bold)
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if vertical == MSO_ANCHOR.MIDDLE:
        tf.paragraphs[0].space_before = Pt(0)

    if not paragraphs:
        return box

    for i, (text, size, color, bold) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        _set_run(run, text, size, color, bold)
    return box


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def accent_bar(slide, left, top, width, height=Inches(0.08)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = EMERALD
    shape.line.fill.background()
    return shape


def card(slide, left, top, width, height, fill: RGBColor = BG_CARD):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    # Soften corners slightly via adjustments if available
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def pill(slide, left, top, width, height, text: str, fill: RGBColor = EMERALD):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, text, 12, WHITE, True)
    return shape


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def numbered_row(slide, left, top, width, number: str, title: str, body: str):
    # Number circle
    circ = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(0.42), Inches(0.42)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = EMERALD
    circ.line.fill.background()
    tf = circ.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, number, 14, BG, True)

    add_textbox(
        slide,
        left + Inches(0.55),
        top - Inches(0.02),
        width - Inches(0.55),
        Inches(0.28),
        paragraphs=[(title, 16, WHITE, True)],
    )
    add_textbox(
        slide,
        left + Inches(0.55),
        top + Inches(0.26),
        width - Inches(0.55),
        Inches(0.35),
        paragraphs=[(body, 12, MUTED, False)],
    )


def flow_chip(slide, left, top, width, height, text: str, accent: bool = False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = EMERALD if accent else BG_CARD
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, text, 11, WHITE if accent else LIGHT, True)
    return shape


def arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, Inches(0.28), Inches(0.18)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = EMERALD
    shape.line.fill.background()
    return shape


# --- Slides ---


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(2.15), Inches(1.6))

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.0),
        Inches(11.5),
        Inches(0.4),
        paragraphs=[("TEAM CJP  ·  HACKATHON PITCH", 14, EMERALD, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.35),
        Inches(11.5),
        Inches(1.0),
        paragraphs=[("Moonshot Feature Hub", 48, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.45),
        Inches(11.5),
        Inches(0.7),
        paragraphs=[
            (
                "Turn every Slack feature ask into revenue-weighted product truth.",
                22,
                MUTED,
                False,
            )
        ],
    )

    # Team members as pills
    members = ["Chidananda", "Sahil", "Akhand", "Roy"]
    x = Inches(0.8)
    for name in members:
        w = Inches(1.55 + len(name) * 0.05)
        pill(slide, x, Inches(5.5), w, Inches(0.42), name, BG_CARD)
        # emerald left edge accent via thin rect
        accent_bar(slide, x, Inches(5.5), Inches(0.08), Inches(0.42))
        x += w + Inches(0.2)

    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.3),
        Inches(11.5),
        Inches(0.35),
        paragraphs=[("Internal leadership & hackathon judges", 13, MUTED, False)],
    )
    add_notes(
        slide,
        "Open strong: introduce Team CJP and the one-liner. "
        "Moonshot Feature Hub turns Slack feature asks into revenue-weighted product decisions. "
        "Members: Chidananda, Sahil, Akhand, Roy.",
    )


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.6),
        paragraphs=[("The Problem", 36, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(11.5),
        Inches(0.4),
        paragraphs=[
            (
                "Customer demand is trapped where Product cannot see it — and Leadership cannot price it.",
                16,
                MUTED,
                False,
            )
        ],
    )

    problems = [
        ("Scattered asks", "Feature requests live in customer Slack threads — invisible to the backlog."),
        ("Tribal CS knowledge", "CS remembers asks; PM systems stay disconnected from that reality."),
        ("No ARR rollup", "The same ask from 5 accounts looks like 5 tickets. Demand is never summed."),
        ("Leadership blind spot", "What demand is unplanned — and how much revenue sits behind it?"),
    ]
    for i, (title, body) in enumerate(problems):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + col * Inches(6.1)
        top = Inches(2.2) + row * Inches(2.1)
        card(slide, left, top, Inches(5.8), Inches(1.85))
        add_textbox(
            slide,
            left + Inches(0.35),
            top + Inches(0.35),
            Inches(5.1),
            Inches(0.4),
            paragraphs=[(title, 20, EMERALD, True)],
        )
        add_textbox(
            slide,
            left + Inches(0.35),
            top + Inches(0.85),
            Inches(5.1),
            Inches(0.75),
            paragraphs=[(body, 14, LIGHT, False)],
        )
    add_notes(
        slide,
        "Paint the pain: Slack is where customers ask, but Product and Leadership cannot "
        "see consolidated, ARR-weighted demand. Same feature from five accounts looks like noise.",
    )


def slide_cost(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.6),
        paragraphs=[("Cost of Inaction", 36, WHITE, True)],
    )

    items = [
        ("Wrong roadmap bets", "Blind spots steer engineering toward loud voices, not revenue."),
        ("Slower ship cycles", "Duplicate and lost signals force re-discovery every quarter."),
        ("Renewal risk", "Open critical asks + dark accounts compound into churn."),
    ]
    for i, (title, body) in enumerate(items):
        left = Inches(0.8) + i * Inches(4.05)
        card(slide, left, Inches(1.9), Inches(3.85), Inches(3.2))
        add_textbox(
            slide,
            left + Inches(0.3),
            Inches(2.2),
            Inches(3.25),
            Inches(0.4),
            paragraphs=[(f"0{i+1}", 28, EMERALD, True)],
        )
        add_textbox(
            slide,
            left + Inches(0.3),
            Inches(3.0),
            Inches(3.25),
            Inches(0.6),
            paragraphs=[(title, 20, WHITE, True)],
        )
        add_textbox(
            slide,
            left + Inches(0.3),
            Inches(3.7),
            Inches(3.25),
            Inches(1.0),
            paragraphs=[(body, 14, MUTED, False)],
        )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11.7),
        Inches(0.8),
        paragraphs=[
            (
                "Demand without ARR is noise.  ARR without demand signal is risk.",
                22,
                WHITE,
                True,
            )
        ],
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "For leadership judges: the cost is misallocated roadmap, slower delivery, and "
        "renewal risk. End on the tagline — demand without ARR is noise; ARR without demand is risk.",
    )


def slide_usp(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.5),
        paragraphs=[("Our USP", 36, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(11.7),
        Inches(1.4),
        paragraphs=[
            (
                "AI captures Slack asks → CS consolidates by theme with ARR → "
                "PM gets a canonical, revenue-weighted backlog → "
                "Retention ties open asks to churn risk.",
                24,
                WHITE,
                True,
            )
        ],
    )

    card(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(2.6))
    add_textbox(
        slide,
        Inches(1.15),
        Inches(3.7),
        Inches(11.0),
        Inches(0.4),
        paragraphs=[("Not another feedback form", 18, EMERALD, True)],
    )
    add_textbox(
        slide,
        Inches(1.15),
        Inches(4.25),
        Inches(11.0),
        Inches(1.4),
        paragraphs=[
            (
                "A closed loop from Slack → triage → consolidation → roadmap → renewal. "
                "Every whisper becomes a decision Product can defend with dollars.",
                18,
                LIGHT,
                False,
            )
        ],
    )
    add_notes(
        slide,
        "State the USP once, slowly. Emphasize closed loop — capture through retention — "
        "not a disconnected intake form.",
    )


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.5),
        paragraphs=[("Solution Snapshot", 36, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(11.5),
        Inches(0.4),
        paragraphs=[
            ("Two tracks that converge on one canonical backlog.", 16, MUTED, False)
        ],
    )

    # CS track card
    card(slide, Inches(0.8), Inches(2.0), Inches(5.7), Inches(3.8))
    add_textbox(
        slide,
        Inches(1.1),
        Inches(2.25),
        Inches(5.1),
        Inches(0.4),
        paragraphs=[("CS / CLM TRACK", 14, EMERALD, True)],
    )
    steps_cs = [
        "1. Product Requests (per workspace)",
        "2. Consolidation themes",
        "3. Unique ARR rollup",
        "4. Promote to PM backlog",
    ]
    for i, s in enumerate(steps_cs):
        add_textbox(
            slide,
            Inches(1.1),
            Inches(2.85) + i * Inches(0.55),
            Inches(5.1),
            Inches(0.45),
            paragraphs=[(s, 16, LIGHT, False)],
        )

    # PM track card
    card(slide, Inches(6.85), Inches(2.0), Inches(5.7), Inches(3.8))
    add_textbox(
        slide,
        Inches(7.15),
        Inches(2.25),
        Inches(5.1),
        Inches(0.4),
        paragraphs=[("PM TRACK", 14, EMERALD, True)],
    )
    steps_pm = [
        "1. Slack Feature Signals",
        "2. Inbox triage (match / create)",
        "3. Canonical Feature Request",
        "4. Roadmap themes & quarters",
    ]
    for i, s in enumerate(steps_pm):
        add_textbox(
            slide,
            Inches(7.15),
            Inches(2.85) + i * Inches(0.55),
            Inches(5.1),
            Inches(0.45),
            paragraphs=[(s, 16, LIGHT, False)],
        )

    # Convergence label
    pill(
        slide,
        Inches(4.15),
        Inches(6.15),
        Inches(5.0),
        Inches(0.45),
        "→  Canonical Feature Request  ←",
        EMERALD_DIM,
    )
    add_notes(
        slide,
        "Explain the two-lane model: CS owns workspace asks and ARR consolidation; "
        "PM owns Slack triage and roadmap. They meet at the canonical Feature Request.",
    )


def slide_flows(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.45), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.65),
        Inches(11.5),
        Inches(0.5),
        paragraphs=[("Major Flows — Full Product Coverage", 32, WHITE, True)],
    )

    flows = [
        ("1", "Capture", "Slack AI intake + CS manual requests + suggestions bell"),
        ("2", "Triage", "Match / Create / Dismiss with similarity ranking"),
        ("3", "Consolidate", "Theme buckets, unique ARR rollup across workspaces"),
        ("4", "Promote", "Consolidation → PM backlog bridge"),
        ("5", "Prioritize", "Analytics: gap board, opportunity prize, revenue at risk"),
        ("6", "Plan", "Roadmap themes/quarters + tags"),
        ("7", "Retain", "Expiry, dark accounts, Slack alerts / nudges"),
    ]
    # Left column 1-4, right column 5-7
    for i, (num, title, body) in enumerate(flows):
        if i < 4:
            left = Inches(0.8)
            top = Inches(1.4) + i * Inches(1.25)
        else:
            left = Inches(7.0)
            top = Inches(1.4) + (i - 4) * Inches(1.25)
        card(slide, left, top, Inches(5.5), Inches(1.1))
        numbered_row(slide, left + Inches(0.2), top + Inches(0.28), Inches(5.1), num, title, body)

    add_notes(
        slide,
        "Walk the seven pillars quickly — Capture, Triage, Consolidate, Promote, "
        "Prioritize, Plan, Retain. This is the full product surface judges should remember.",
    )


def slide_journey(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.5),
        paragraphs=[("End-to-End Journey", 36, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(11.5),
        Inches(0.4),
        paragraphs=[
            ("One linear path we will demo live — from whisper to shipped.", 16, MUTED, False)
        ],
    )

    # Row 1
    row1 = [
        "Slack\nmessage",
        "AI\nsignal",
        "Inbox\ntriage",
        "Match or\nCreate",
    ]
    row2 = [
        "Consolidation\nARR",
        "Analytics\ngap / prize",
        "Promote /\nRoadmap",
        "Retain →\nSHIPPED",
    ]

    chip_w = Inches(2.4)
    chip_h = Inches(1.0)
    gap = Inches(0.45)
    start_x = Inches(0.9)

    for i, label in enumerate(row1):
        x = start_x + i * (chip_w + gap)
        flow_chip(slide, x, Inches(2.3), chip_w, chip_h, label, accent=(i == 0))
        if i < len(row1) - 1:
            arrow_right(slide, x + chip_w + Inches(0.08), Inches(2.7))

    for i, label in enumerate(row2):
        x = start_x + i * (chip_w + gap)
        accent = i == len(row2) - 1
        flow_chip(slide, x, Inches(4.0), chip_w, chip_h, label, accent=accent)
        if i < len(row2) - 1:
            arrow_right(slide, x + chip_w + Inches(0.08), Inches(4.4))

    # Down arrow hint between rows
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.4),
        Inches(11.7),
        Inches(0.4),
        paragraphs=[("↓", 20, EMERALD, True)],
        align=PP_ALIGN.CENTER,
    )

    card(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2))
    add_textbox(
        slide,
        Inches(1.15),
        Inches(5.8),
        Inches(11.0),
        Inches(0.6),
        paragraphs=[
            (
                "We’ll demo this live next.  Backup screen recording ready if wifi fails.",
                18,
                WHITE,
                True,
            )
        ],
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Hand off to live demo here. Follow this exact path. "
        "Have a 60-second screen capture ready as wifi insurance.",
    )


def slide_built(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.5),
        paragraphs=[("What’s Built", 36, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(11.5),
        Inches(0.4),
        paragraphs=[("Proof of execution — not a mock.", 16, MUTED, False)],
    )

    modules = [
        "Inbox",
        "Triage",
        "Requests",
        "Backlog",
        "Consolidation",
        "Roadmap",
        "Analytics",
        "Retention",
        "Orgs",
        "Tags",
        "CS Owners",
        "Settings",
    ]
    for i, name in enumerate(modules):
        col = i % 4
        row = i // 4
        left = Inches(0.8) + col * Inches(3.05)
        top = Inches(2.0) + row * Inches(1.0)
        card(slide, left, top, Inches(2.9), Inches(0.8))
        add_textbox(
            slide,
            left,
            top + Inches(0.2),
            Inches(2.9),
            Inches(0.45),
            paragraphs=[(name, 16, WHITE, True)],
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.4),
        Inches(11.7),
        Inches(0.8),
        paragraphs=[
            (
                "Stack:  Next.js 15  ·  Postgres / Prisma  ·  Slack Events  ·  OpenAI classifier  ·  NextAuth",
                16,
                MUTED,
                False,
            )
        ],
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Feasibility slide for judges: this is a working hub across twelve surfaces. "
        "Stack is production-shaped — Next.js, Postgres, Slack, OpenAI classifier.",
    )


def slide_impact(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.5),
        paragraphs=[("Impact for the Business", 36, WHITE, True)],
    )

    impacts = [
        ("One source of truth", "CS and Product share the same demand record — no more parallel spreadsheets."),
        ("Prize score", "Accounts × unique ARR. Prioritize what moves revenue, not what’s loudest."),
        ("Gap board", "Multi-customer demand not yet on the roadmap — the prize sitting unplanned."),
        ("Retention radar", "Contract expiry + usage decline + open critical asks in one view."),
    ]
    for i, (title, body) in enumerate(impacts):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + col * Inches(6.1)
        top = Inches(1.7) + row * Inches(2.0)
        card(slide, left, top, Inches(5.8), Inches(1.8))
        add_textbox(
            slide,
            left + Inches(0.35),
            top + Inches(0.3),
            Inches(5.1),
            Inches(0.4),
            paragraphs=[(title, 20, EMERALD, True)],
        )
        add_textbox(
            slide,
            left + Inches(0.35),
            top + Inches(0.85),
            Inches(5.1),
            Inches(0.7),
            paragraphs=[(body, 14, LIGHT, False)],
        )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.0),
        Inches(11.7),
        Inches(0.5),
        paragraphs=[
            ("Ship the features that protect and grow ARR.", 20, WHITE, True)
        ],
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Leadership language: source of truth, prize score, gap board, retention radar. "
        "Close with: ship the features that protect and grow ARR.",
    )


def slide_wow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(1.5), Inches(1.6))

    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.9),
        Inches(11.5),
        Inches(0.4),
        paragraphs=[("THE ASK", 14, EMERALD, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.7),
        Inches(11.7),
        Inches(1.2),
        paragraphs=[
            (
                "Adopt Moonshot as the operating system for feature demand.",
                32,
                WHITE,
                True,
            )
        ],
    )

    nexts = [
        "Deepen Jira sync",
        "Expand classifier coverage",
        "Operationalize retention alerts",
    ]
    for i, item in enumerate(nexts):
        left = Inches(0.8) + i * Inches(4.05)
        card(slide, left, Inches(3.3), Inches(3.85), Inches(1.1))
        add_textbox(
            slide,
            left + Inches(0.25),
            Inches(3.6),
            Inches(3.35),
            Inches(0.5),
            paragraphs=[(item, 16, LIGHT, True)],
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.8),
        Inches(11.7),
        Inches(0.9),
        paragraphs=[
            (
                "Every customer whisper in Slack becomes a decision Product can defend with dollars.",
                18,
                MUTED,
                False,
            )
        ],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.9),
        Inches(11.7),
        Inches(0.5),
        paragraphs=[("Team CJP  —  Thank you  ·  Questions", 20, WHITE, True)],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.5),
        Inches(11.7),
        Inches(0.35),
        paragraphs=[("Chidananda  ·  Sahil  ·  Akhand  ·  Roy", 14, MUTED, False)],
        align=PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        "Ask clearly: adopt Moonshot as the OS for feature demand. "
        "Name next steps briefly, land the vision line, then open for questions as Team CJP.",
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.45),
        paragraphs=[("Backup: Architecture", 32, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.3),
        Inches(11.5),
        Inches(0.35),
        paragraphs=[("Q&A only — not part of the timed pitch.", 14, MUTED, False)],
    )

    stages = [
        ("Slack Events", "Message lands in mapped channel"),
        ("Ingest", "handleSlackMessage validates & routes"),
        ("Classifier", "OpenAI / heuristic feature-request score"),
        ("Feature Signal", "PENDING in Inbox with AI title/tags"),
        ("Triage APIs", "match · create · dismiss · similar"),
        ("Canonical FR", "Backlog + roadmap + activity timeline"),
    ]
    for i, (title, body) in enumerate(stages):
        col = i % 3
        row = i // 3
        left = Inches(0.8) + col * Inches(4.05)
        top = Inches(2.0) + row * Inches(2.2)
        card(slide, left, top, Inches(3.85), Inches(1.9))
        add_textbox(
            slide,
            left + Inches(0.3),
            top + Inches(0.35),
            Inches(3.25),
            Inches(0.35),
            paragraphs=[(f"{i+1}. {title}", 16, EMERALD, True)],
        )
        add_textbox(
            slide,
            left + Inches(0.3),
            top + Inches(0.9),
            Inches(3.25),
            Inches(0.7),
            paragraphs=[(body, 14, LIGHT, False)],
        )
    add_notes(
        slide,
        "Backup for technical Q&A: Slack → ingest → classifier → signal → triage APIs → canonical feature request.",
    )


def slide_datamodel(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    accent_bar(slide, Inches(0.8), Inches(0.55), Inches(1.2))
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(11.5),
        Inches(0.45),
        paragraphs=[("Backup: Data Model", 32, WHITE, True)],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.3),
        Inches(11.5),
        Inches(0.35),
        paragraphs=[("Q&A only — core entities that power the hub.", 14, MUTED, False)],
    )

    entities = [
        ("CustomerOrg", "Workspace · ARR · contract · activity"),
        ("FeatureSignal", "AI-detected Slack ask (PENDING/MATCHED)"),
        ("ProductRequest", "CS/CLM ask per workspace"),
        ("Consolidation", "Theme bucket · derived unique ARR"),
        ("FeatureRequest", "Canonical PM backlog item"),
        ("RoadmapItem", "Theme / quarter grouping"),
        ("Suggestion", "Alternate intake with match %"),
        ("RetentionAlert", "Logged alerts & nudges"),
    ]
    for i, (title, body) in enumerate(entities):
        col = i % 4
        row = i // 4
        left = Inches(0.8) + col * Inches(3.05)
        top = Inches(2.0) + row * Inches(2.2)
        card(slide, left, top, Inches(2.9), Inches(1.9))
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(0.4),
            Inches(2.5),
            Inches(0.5),
            paragraphs=[(title, 15, EMERALD, True)],
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            top + Inches(1.0),
            Inches(2.5),
            Inches(0.65),
            paragraphs=[(body, 12, LIGHT, False)],
        )
    add_notes(
        slide,
        "Backup data model: Org, Signal, ProductRequest, Consolidation, FeatureRequest, "
        "RoadmapItem, Suggestion, RetentionAlert.",
    )


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_cost(prs)
    slide_usp(prs)
    slide_solution(prs)
    slide_flows(prs)
    slide_journey(prs)
    slide_built(prs)
    slide_impact(prs)
    slide_wow(prs)
    slide_architecture(prs)
    slide_datamodel(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
